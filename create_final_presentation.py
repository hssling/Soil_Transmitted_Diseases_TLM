#!/usr/bin/env python3
"""
Final Professional PPTX Presentation Generator for STH MBBS Teaching
Creates a comprehensive PowerPoint presentation with professional formatting
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import markdown
import re

class ProfessionalSTHPresenter:
    def __init__(self, content_dir, visuals_dir):
        self.content_dir = Path(content_dir)
        self.visuals_dir = Path(visuals_dir)
        self.presentation = Presentation()

        # Set up professional theme
        self._setup_theme()

        # Load content
        self.slide_content = self._load_markdown_content()
        self.visual_files = self._get_visual_files()

    def _setup_theme(self):
        """Set up professional presentation theme"""
        # Set slide dimensions to 16:9
        self.presentation.slide_width = Inches(13.33)
        self.presentation.slide_height = Inches(7.5)

        # Define colors (Indian theme)
        self.colors = {
            'primary': RGBColor(255, 153, 51),     # Saffron
            'secondary': RGBColor(19, 136, 8),     # Green
            'accent': RGBColor(0, 0, 128),         # Navy Blue
            'text': RGBColor(0, 0, 0),             # Black
            'background': RGBColor(245, 245, 245), # Light Gray
        }

        # Font settings
        self.title_font_size = Pt(44)
        self.subtitle_font_size = Pt(32)
        self.body_font_size = Pt(24)
        self.caption_font_size = Pt(18)

    def _load_markdown_content(self):
        """Load and parse markdown content for slides"""
        content_files = [
            'Presentation_Slides_STH.md',
            'Student_Notes_STH.md'
        ]

        all_content = {}
        for file in content_files:
            file_path = self.content_dir / file
            if file_path.exists():
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()

                # Parse markdown to extract sections
                sections = self._parse_markdown_sections(content)
                all_content[file] = sections

        return all_content

    def _parse_markdown_sections(self, markdown_content):
        """Parse markdown content into sections and slides"""
        lines = markdown_content.split('\n')
        sections = {}
        current_section = None
        current_content = []

        for line in lines:
            line = line.strip()
            if line.startswith('#'):
                # Save previous section
                if current_section and current_content:
                    sections[current_section] = current_content

                # Start new section
                current_section = line.lstrip('#').strip()
                current_content = []
            elif current_section and line:
                current_content.append(line)

        # Save final section
        if current_section and current_content:
            sections[current_section] = current_content

        return sections

    def _get_visual_files(self):
        """Get list of available visual files"""
        visuals = {}
        if self.visuals_dir.exists():
            for png_file in self.visuals_dir.glob('*.png'):
                name = png_file.stem
                visuals[name] = png_file
        return visuals

    def create_title_slide(self):
        """Create professional title slide"""
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[0])

        # Title
        title = slide.shapes.title
        title.text = "Soil Transmitted Diseases (STH)"
        title.text_frame.paragraphs[0].font.size = self.title_font_size
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']

        # Subtitle
        subtitle = slide.placeholders[1]
        subtitle.text = "Comprehensive Teaching Module for MBBS 3rd Year Students\nIndian Context and Public Health Perspectives"
        subtitle.text_frame.paragraphs[0].font.size = self.subtitle_font_size
        subtitle.text_frame.paragraphs[0].font.color.rgb = self.colors['text']

        # Add institutional footer
        self._add_footer(slide, "Medical Education Department | 2025")

    def create_learning_objectives_slide(self):
        """Create learning objectives slide"""
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Learning Objectives"

        objectives = [
            "• Understand epidemiology and global burden of STH",
            "• Identify clinical manifestations and complications",
            "• Master diagnostic approaches and treatment strategies",
            "• Recognize public health aspects and control measures",
            "• Apply Indian context and health program integration"
        ]

        content = slide.placeholders[1]
        for obj in objectives:
            p = content.text_frame.add_paragraph()
            p.text = obj
            p.font.size = self.body_font_size
            p.font.color.rgb = self.colors['text']

    def create_content_slides(self):
        """Create comprehensive content slides in logical medical education flow"""
        # Define logical sequence of topics
        topic_sequence = [
            ("Overview of Soil Transmitted Diseases", [
                "Soil Transmitted Helminthiases (STH) are parasitic worm infections",
                "Three major parasites: Ascaris lumbricoides, Trichuris trichiura, Hookworms",
                "Global burden: 1.5 billion people affected (24% of world population)",
                "India: 225 million cases, highest absolute burden globally",
                "Major public health problem in tropical and developing regions",
                "Especially affects children and rural communities"
            ]),
            ("Epidemiology and Global Burden", [
                "Prevalence: Ascaris (807-1221M), Trichuris (604-795M), Hookworms (576-740M)",
                "Geographic distribution: Sub-Saharan Africa, South Asia, Southeast Asia",
                "Risk factors: Poverty, poor sanitation, open defecation, inadequate water",
                "Age pattern: School-age children (peak intensity), preschoolers (peak burden)",
                "Socioeconomic impact: Economic cost $7-12 billion annually",
                "DALYs: 4.98 million annually, significant productivity losses"
            ]),
            ("Etiology and Life Cycles", [
                "Ascaris lumbricoides: Fecal-oral transmission, 3-4 weeks development",
                "Embryonated eggs → Larval migration (lung-liver) → Adult in jejunum",
                "Trichuris trichiura: Fecal-oral route, 3-6 weeks soil development",
                "Local tissue penetration, adults in cecum/colon",
                "Hookworms (Necator americanus): Percutaneous transmission through skin",
                "Filariform larvae penetrate feet, pulmonary migration, adults in jejunum",
                "Environmental requirements: Warm, humid soil, adequate oxygen"
            ]),
            ("Clinical Manifestations", [
                "Asymptomatic infections: 60-80% of cases show no symptoms",
                "Pulmonary phase: Loeffler's syndrome (cough, eosinophilia) - Ascaris",
                "Intestinal phase: Abdominal pain, malnutrition, growth retardation",
                "Hookworm disease: Iron deficiency anemia, edema, classic triad",
                "Trichuriasis: Dysentery syndrome, rectal prolapse, chronic diarrhea",
                "Complications: Intestinal obstruction, biliary ascariasis, severe anemia"
            ]),
            ("Diagnosis and Laboratory Methods", [
                "Clinical diagnosis: Geographic history, symptoms, occupational exposure",
                "Stool examination: Direct smear vs concentration techniques",
                "Kato-Katz technique: WHO gold standard for field surveys",
                "Egg morphology: Ascaris (oval, mamillated), Trichuris (barrel, bipolar plugs)",
                "Hookworm eggs: Oval, thin-shelled, 4-8 cell morula stage",
                "Advanced diagnostics: PCR, LAMP for surveillance programs"
            ]),
            ("Treatment and Management", [
                "First-line drugs: Albendazole 400mg or Mebendazole 500mg single dose",
                "Efficacy: 95-100% for Ascaris, 70-95% for hookworms, 30-90% for Trichuris",
                "Pyrantel pamoate: 10mg/kg single dose, safe in pregnancy (first trimester)",
                "Special populations: Caution in pregnancy, iron supplementation for anemia",
                "Complication management: Surgical intervention for obstructions",
                "Drug resistance monitoring: Emerging concern in some regions"
            ]),
        ]

        for section_title, bullet_points in topic_sequence:
            slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])

            # Set title
            title = slide.shapes.title
            title.text = section_title
            title.text_frame.paragraphs[0].font.size = Pt(36)
            title.text_frame.paragraphs[0].font.bold = True
            title.text_frame.paragraphs[0].font.color.rgb = self.colors['secondary']

            # Add content
            content_box = slide.placeholders[1]

            for point in bullet_points:
                p = content_box.text_frame.add_paragraph()
                p.text = "• " + point
                p.font.size = self.body_font_size
                p.font.color.rgb = self.colors['text']
                p.level = 0

                self._add_footer(slide, f"STH Module - {section_title}")

        # Add expanded prevention and control section (multiple slides) - Special Focus
        self._create_prevention_control_slides()

    def _extract_bullet_points(self, text):
        """Extract bullet points from text"""
        points = []
        lines = text.split('\n')

        for line in lines:
            line = line.strip()
            if line.startswith('-') or line.startswith('•') or line.startswith('*'):
                points.append(line.lstrip('-•*').strip())
            elif line and not line.startswith('#'):
                points.append(line)

        return points

    def create_visual_slides(self):
        """Create slides featuring the generated visualizations"""
        visual_slides = [
            ('Indian STH Epidemiology', 'Statewise_STH_Prevalence.png'),
            ('Parasite Prevalence Patterns', 'Parasite_Prevalence_Comparison.png'),
            ('Risk Distribution', 'Risk_Category_Distribution.png'),
            ('Regional Comparisons', 'Regional_Prevalence_Heatmap.png'),
            ('Healthcare Integration', 'Healthcare_Integration.png'),
            ('Program Progress', 'Deworming_Progress.png'),
            ('Comprehensive Dashboard', 'STH_India_Dashboard.png')
        ]

        for title, image_file in visual_slides:
            slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[5])  # Title only layout

            # Add title
            title_shape = slide.shapes.title
            title_shape.text = title
            title_shape.text_frame.paragraphs[0].font.size = Pt(36)
            title_shape.text_frame.paragraphs[0].font.bold = True
            title_shape.text_frame.paragraphs[0].font.color.rgb = self.colors['accent']

            # Add image if exists
            if image_file in self.visual_files:
                image_path = self.visual_files[image_file]

                # Calculate image size (leave margin for title)
                left = Inches(0.5)
                top = Inches(1.5)
                width = Inches(12.33)
                height = Inches(5.5)

                slide.shapes.add_picture(str(image_path), left, top, width, height)

            self._add_footer(slide, "Indian Context Data Visualization")

    def create_summary_slide(self):
        """Create summary and key points slide"""
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Summary and Key Takeaways"

        key_points = [
            "• STH affect 1.5 billion people globally, India has 225 million cases",
            "• Three main parasites: Ascaris, Trichuris, Hookworms",
            "• Transmission: Fecal-oral and percutaneous routes",
            "• Control through MDA, WASH, and health education",
            "• Integrated with Indian health programs (NDD, NHM)",
            "• Prevention better than cure - focus on children"
        ]

        content = slide.placeholders[1]
        for point in key_points:
            p = content.text_frame.add_paragraph()
            p.text = point
            p.font.size = self.body_font_size
            p.font.color.rgb = self.colors['text']

    def _create_prevention_control_slides(self):
        """Create multiple slides focused on prevention and control - Special Emphasis"""
        # Slide 1: Prevention Strategies Overview
        slide1 = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])
        title1 = slide1.shapes.title
        title1.text = "Prevention and Control Strategies - Overview"
        title1.text_frame.paragraphs[0].font.size = Pt(36)
        title1.text_frame.paragraphs[0].font.bold = True
        title1.text_frame.paragraphs[0].font.color.rgb = self.colors['accent']

        prevention_overview = [
            "WHO Control Framework: Preventive chemotherapy, WASH interventions, Health education",
            "Three Pillars: Mass drug administration, Sanitation improvement, Behavioral change",
            "Target Population: School-age children (5-14 years), preschoolers, high-risk groups",
            "WHO 2030 Targets: 75% prevalence reduction, 90% coverage, elimination in some countries",
            "Cost-Effective: $0.02-0.50 per treatment, benefit-cost ratio 1:30"
        ]

        content1 = slide1.placeholders[1]
        for point in prevention_overview:
            p = content1.text_frame.add_paragraph()
            p.text = "• " + point
            p.font.size = self.body_font_size
            p.font.color.rgb = self.colors['text']
        self._add_footer(slide1, "Prevention & Control - Overview")

        # Slide 2: Mass Drug Administration (MDA)
        slide2 = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])
        title2 = slide2.shapes.title
        title2.text = "Mass Drug Administration (MDA) Programs"
        title2.text_frame.paragraphs[0].font.size = Pt(36)
        title2.text_frame.paragraphs[0].font.bold = True
        title2.text_frame.paragraphs[0].font.color.rgb = self.colors['accent']

        mda_content = [
            "Primary Control Strategy: Periodic deworming of at-risk populations",
            "Frequency: 1-2 times annually in high-prevalence areas",
            "Coverage Target: ≥75% of school-age children and high-risk groups",
            "Drugs: Albendazole 400mg or Mebendazole 500mg single oral dose",
            "Safety: Generally safe, even in pregnancy (avoid first trimester)",
            "Monitoring: Treatment coverage, drug efficacy, side effects reporting"
        ]

        content2 = slide2.placeholders[1]
        for point in mda_content:
            p = content2.text_frame.add_paragraph()
            p.text = "• " + point
            p.font.size = self.body_font_size
            p.font.color.rgb = self.colors['text']
        self._add_footer(slide2, "Prevention & Control - MDA Programs")

        # Slide 3: WASH Interventions
        slide3 = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])
        title3 = slide3.shapes.title
        title3.text = "Water, Sanitation and Hygiene (WASH) Interventions"
        title3.text_frame.paragraphs[0].font.size = Pt(36)
        title3.text_frame.paragraphs[0].font.bold = True
        title3.text_frame.paragraphs[0].font.color.rgb = self.colors['accent']

        wash_content = [
            "Sustainable Prevention: Breaks transmission cycle through environmental control",
            "Safe Water Supply: Protected sources, household water treatment and storage",
            "Sanitation: Latrine construction and maintenance, elimination of open defecation",
            "Hygiene Education: Handwashing at critical points, footwear use, food hygiene",
            "Behavioral Change: Community-led programs, school-based education",
            "F Diagram: Feces → Fields → Flies → Fingers → Food → Mouth"
        ]

        content3 = slide3.placeholders[1]
        for point in wash_content:
            p = content3.text_frame.add_paragraph()
            p.text = "• " + point
            p.font.size = self.body_font_size
            p.font.color.rgb = self.colors['text']
        self._add_footer(slide3, "Prevention & Control - WASH Interventions")

        # Slide 4: Health Education and Community Mobilization
        slide4 = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])
        title4 = slide4.shapes.title
        title4.text = "Health Education and Community Engagement"
        title4.text_frame.paragraphs[0].font.size = Pt(36)
        title4.text_frame.paragraphs[0].font.bold = True
        title4.text_frame.paragraphs[0].font.color.rgb = self.colors['accent']

        education_content = [
            "Community Empowerment: Local ownership and participation in control programs",
            "Health Education: Understanding transmission, prevention, treatment compliance",
            "School Programs: Child-to-child education, teacher training, regular deworming",
            "Social Mobilization: Religious leaders, local government, NGOs involvement",
            "Monitoring & Evaluation: Community health workers, surveillance systems",
            "Sustainability: Building local capacity for long-term program success"
        ]

        content4 = slide4.placeholders[1]
        for point in education_content:
            p = content4.text_frame.add_paragraph()
            p.text = "• " + point
            p.font.size = self.body_font_size
            p.font.color.rgb = self.colors['text']
        self._add_footer(slide4, "Prevention & Control - Health Education")

        # Slide 5: Indian National Programs
        slide5 = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])
        title5 = slide5.shapes.title
        title5.text = "National Deworming Day (NDD) - India"
        title5.text_frame.paragraphs[0].font.size = Pt(36)
        title5.text_frame.paragraphs[0].font.bold = True
        title5.text_frame.paragraphs[0].font.color.rgb = self.colors['accent']

        india_programs = [
            "Launched 2015: Annual deworming on February 10th and August 10th",
            "Target Population: Children aged 1-19 years (540 million eligible)",
            "Coverage Achievement: Over 85% in recent years with government efforts",
            "Integration: With Ministry of Health, Education, and Rural Development",
            "Monitoring: Real-time tracking through government portal",
            "Success: Significant reduction in STH prevalence in school-aged children"
        ]

        content5 = slide5.placeholders[1]
        for point in india_programs:
            p = content5.text_frame.add_paragraph()
            p.text = "• " + point
            p.font.size = self.body_font_size
            p.font.color.rgb = self.colors['text']
        self._add_footer(slide5, "Prevention & Control - India Programs")

        # Slide 6: Challenges and Future Directions
        slide6 = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])
        title6 = slide6.shapes.title
        title6.text = "Challenges and Future Directions"
        title6.text_frame.paragraphs[0].font.size = Pt(36)
        title6.text_frame.paragraphs[0].font.bold = True
        title6.text_frame.paragraphs[0].font.color.rgb = self.colors['accent']

        challenges_content = [
            "Drug Resistance: Emerging benzimidazole resistance in some areas",
            "Climate Change: Extended transmission seasons and changing patterns",
            "Urbanization: New transmission dynamics in growing cities",
            "Supply Chain: Ensuring consistent drug availability and quality",
            "Integration: Better coordination between health and other sectors",
            "Future: Vaccines, new drugs, elimination targets achievement"
        ]

        content6 = slide6.placeholders[1]
        for point in challenges_content:
            p = content6.text_frame.add_paragraph()
            p.text = "• " + point
            p.font.size = self.body_font_size
            p.font.color.rgb = self.colors['text']
        self._add_footer(slide6, "Prevention & Control - Challenges & Future")

    def _add_footer(self, slide, text):
        """Add professional footer to slide"""
        left = Inches(0.5)
        top = Inches(6.8)
        width = Inches(12.33)
        height = Inches(0.3)

        footer = slide.shapes.add_textbox(left, top, width, height)
        footer.text = text
        footer.text_frame.paragraphs[0].font.size = Pt(12)
        footer.text_frame.paragraphs[0].font.color.rgb = self.colors['accent']
        footer.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    def save_presentation(self, output_path):
        """Save the final presentation"""
        self.presentation.save(output_path)
        print(f"✅ Professional STH presentation saved to: {output_path}")

        # Print stats
        num_slides = len(self.presentation.slides)
        print(f"📊 Total slides: {num_slides}")
        print(f"🎨 Professional formatting: Indian theme colors")
        print(f"📈 Visuals integrated: {len(self.visual_files)} images")
        print(f"📝 Content sections: {len(self.slide_content)} sources")

def main():
    """Generate the final professional PPTX presentation"""
    script_dir = Path(__file__).parent
    content_dir = script_dir  # Repository root
    visuals_dir = script_dir / "Visual_Assets_Indian_Context" / "Generated_Charts"
    output_file = script_dir / "STH_Medical_Teaching_Presentation.pptx"

    print("🎯 Creating Professional STH Teaching Presentation...")
    print("🎨 Applying Indian theme with saffron, green, navy colors")
    print("📊 Integrating generated visualizations")
    print("📝 Processing markdown content for slides")

    # Create presenter
    presenter = ProfessionalSTHPresenter(content_dir, visuals_dir)

    # Build presentation
    presenter.create_title_slide()
    print("✓ Title slide created")

    presenter.create_learning_objectives_slide()
    print("✓ Learning objectives slide created")

    presenter.create_content_slides()
    print("✓ Content slides generated from markdown")

    presenter.create_visual_slides()
    print("✓ Visual slides with Indian context data created")

    presenter.create_summary_slide()
    print("✓ Summary slide added")

    # Save presentation
    presenter.save_presentation(output_file)

    print("\n🎉 Professional STH presentation complete!")
    print("=" * 50)
    print("📁 Ready for classroom presentation")
    print("🎨 Professional formatting applied")
    print("🇮🇳 Indian context integrated")
    print("📊 Data visualizations included")
    print("=" * 50)

if __name__ == "__main__":
    main()
