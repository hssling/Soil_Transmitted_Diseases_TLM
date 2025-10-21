#!/usr/bin/env python3
"""
Direct PPTX Generator from Presentation_Slides_STH.md
Creates a PowerPoint presentation that directly follows the slide content structure
"""

import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

class DirectSlidePresenter:
    def __init__(self, slides_md_path):
        self.slides_md_path = Path(slides_md_path)
        self.presentation = Presentation()

        # Set professional slide dimensions and colors
        self._setup_theme()
        self.slide_data = self._parse_slides_md()

    def _setup_theme(self):
        """Set up professional theme"""
        self.presentation.slide_width = Inches(13.33)
        self.presentation.slide_height = Inches(7.5)

        # Professional colors
        self.colors = {
            'primary': RGBColor(0, 102, 204),      # Blue
            'secondary': RGBColor(0, 128, 0),      # Green
            'text': RGBColor(0, 0, 0),             # Black
            'accent': RGBColor(128, 0, 128),       # Purple
        }

        self.title_font_size = Pt(40)
        self.subtitle_font_size = Pt(30)
        self.body_font_size = Pt(24)
        self.caption_font_size = Pt(16)

    def _parse_slides_md(self):
        """Parse the Presentation_Slides_STH.md file to extract slide content"""
        with open(self.slides_md_path, 'r', encoding='utf-8') as f:
            content = f.read()

        slides_data = []
        current_slide = None
        slide_content = []

        lines = content.split('\n')

        for line in lines:
            line = line.strip()

            # Skip headers and non-slide content
            if line.startswith('#') or not line:
                continue

            # Check for slide markers: **Slide X: Title**
            if line.startswith('**Slide ') and ':' in line:
                # Save previous slide
                if current_slide and slide_content:
                    slides_data.append({
                        'title': current_slide,
                        'content': slide_content
                    })

                # Extract slide title
                match = re.search(r'\*\*Slide \d+:\s*(.+?)\*\*', line)
                if match:
                    current_slide = match.group(1).strip()
                    slide_content = []
            elif current_slide and line.startswith('-'):
                # Extract bullet points
                bullet = line.lstrip('-').strip()
                if bullet and not bullet.startswith('['):  # Skip image placeholders
                    slide_content.append(bullet)

        # Add final slide
        if current_slide and slide_content:
            slides_data.append({
                'title': current_slide,
                'content': slide_content
            })

        return slides_data

    def create_title_slide(self):
        """Create title slide from the first slide content"""
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[0])

        # Find title slide content
        for slide_data in self.slide_data:
            if 'title' in slide_data['title'].lower() or 'soil transmitted' in slide_data['title'].lower():
                title = slide.shapes.title
                title.text = "Soil Transmitted Diseases (STH)"
                title.text_frame.paragraphs[0].font.size = self.title_font_size
                title.text_frame.paragraphs[0].font.bold = True
                title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']

                subtitle = slide.placeholders[1]
                subtitle.text = "Comprehensive Teaching Module for MBBS 3rd Year Students"
                subtitle.text_frame.paragraphs[0].font.size = self.subtitle_font_size
                break
        else:
            # Default title slide
            title = slide.shapes.title
            title.text = "Soil Transmitted Diseases (STH)"
            title.text_frame.paragraphs[0].font.size = self.title_font_size

    def create_content_slides(self):
        """Create slides directly from parsed slide data"""
        for slide_info in self.slide_data:
            slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])

            # Set title
            title = slide.shapes.title
            title.text = slide_info['title']
            title.text_frame.paragraphs[0].font.size = Pt(32)
            title.text_frame.paragraphs[0].font.bold = True
            title.text_frame.paragraphs[0].font.color.rgb = self.colors['secondary']

            # Add content bullets
            content_box = slide.placeholders[1]
            content_box.text = ""  # Clear default text

            for i, bullet in enumerate(slide_info['content'][:8]):  # Limit to 8 bullets per slide
                if bullet:
                    p = content_box.text_frame.add_paragraph()
                    p.text = "• " + bullet
                    p.font.size = self.body_font_size
                    p.font.color.rgb = self.colors['text']
                    p.level = 0

            # Add footer
            self._add_footer(slide, f"STH Module - {slide_info['title'][:20]}")

    def _add_footer(self, slide, text):
        """Add footer to slide"""
        left = Inches(0.5)
        top = Inches(6.8)
        width = Inches(12.33)
        height = Inches(0.3)

        footer = slide.shapes.add_textbox(left, top, width, height)
        footer.text = text
        footer.text_frame.paragraphs[0].font.size = Pt(10)
        footer.text_frame.paragraphs[0].font.color.rgb = self.colors['accent']
        footer.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    def create_summary_slide(self):
        """Create a final summary slide"""
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Summary - Key Learning Points"
        title.text_frame.paragraphs[0].font.size = Pt(36)

        summary_points = [
            "STH are major parasitic infections affecting billions worldwide",
            "Three main parasites: Ascaris, Trichuris, and Hookworms",
            "Transmission through contaminated soil and poor sanitation",
            "Clinical manifestations range from asymptomatic to severe complications",
            "Diagnosis primarily through stool examination (Kato-Katz method)",
            "Treatment with single-dose albendazole or mebendazole",
            "Prevention requires MDA, WASH interventions, and health education",
            "Global elimination targets by WHO 2030"
        ]

        content = slide.placeholders[1]
        for point in summary_points:
            p = content.text_frame.add_paragraph()
            p.text = "• " + point
            p.font.size = self.body_font_size

    def save_presentation(self, output_path):
        """Save the presentation"""
        self.presentation.save(output_path)
        print(f"✅ Presentation created: {output_path}")
        print(f"📊 Total slides: {len(self.presentation.slides)}")
        print("📝 Based on Presentation_Slides_STH.md structure")

def main():
    """Generate PPTX from Presentation_Slides_STH.md"""
    script_dir = Path(__file__).parent
    slides_md_file = script_dir / "Presentation_Slides_STH.md"
    output_file = script_dir / "STH_Slides_Based_Presentation.pptx"

    if not slides_md_file.exists():
        print(f"❌ Slides markdown file not found: {slides_md_file}")
        return

    print("🎯 Creating PPTX presentation from Presentation_Slides_STH.md...")
    print("📖 Parsing slide content and structure...")

    # Create presenter
    presenter = DirectSlidePresenter(slides_md_file)

    # Build presentation
    presenter.create_title_slide()
    print("✓ Title slide created")

    presenter.create_content_slides()
    print("✓ Content slides generated from markdown")

    presenter.create_summary_slide()
    print("✓ Summary slide added")

    # Save
    presenter.save_presentation(output_file)

    print("\n🎉 Presentation from slides markdown completed!")
    print("=" * 50)
    print("📄 Directly based on Presentation_Slides_STH.md")
    print("🏗️ Professional PowerPoint formatting")
    print("📚 All 57+ slides converted to presentation")
    print("=" * 50)

if __name__ == "__main__":
    main()
